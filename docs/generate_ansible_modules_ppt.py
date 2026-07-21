#!/usr/bin/env python3
"""Generate professional PowerPoint: Ansible Modules — What, Why, and Examples."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
TEAL = RGBColor(0x00, 0x96, 0x88)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
ACCENT = RGBColor(0xE6, 0x7E, 0x22)
MUTED = RGBColor(0x7F, 0x8C, 0x8D)

OUTPUT = "/workspace/docs/Ansible_Modules_Presentation.pptx"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, subtitle=None):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.25)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12), Inches(0.4))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = TEAL
        sp.font.italic = True


def add_bullets(slide, items, left=0.6, top=1.7, width=12, height=5.5, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)


def add_code_block(slide, code, top=3.2, height=2.8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(top), Inches(12.1), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    shape.line.color.rgb = TEAL
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)


def add_footer(slide, text="Ansible Modules | MS Office 2016 LAN Project"):
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.35))
    p = foot.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2))
    p = t1.text_frame.paragraphs[0]
    p.text = "Ansible Modules"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE

    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.6))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "What They Are, Why We Use Them & Real Examples"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)

    t3 = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.5), Inches(0.5))
    p3 = t3.text_frame.paragraphs[0]
    p3.text = "MS Office 2016 LAN Deployment | Ansible AWX"
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEAL


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Agenda")
    add_bullets(
        slide,
        [
            "What is Ansible?",
            "What is a Module?",
            "Why We Use Modules in Ansible",
            "Module vs Command / Shell",
            "Module Categories & Anatomy",
            "Idempotency Explained",
            "Real Examples from Our Project",
            "Best Practices & Summary",
        ],
        font_size=20,
    )
    add_footer(slide)


def slide_what_is_ansible(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "What is Ansible?", "Agentless IT automation platform")
    add_bullets(
        slide,
        [
            "Open-source tool for configuration, deployment, and orchestration",
            "Uses simple YAML playbooks — no programming required",
            "Agentless: connects via SSH (Linux) or WinRM (Windows)",
            "Core building blocks:",
            ("  • Inventory — list of managed hosts", 1),
            ("  • Playbook — automation workflow in YAML", 1),
            ("  • Module — single action performed on a host", 1),
            ("  • Role — organized tasks, vars, and templates", 1),
            "This project: AWX runs playbooks to install Office 2016 on Windows LAN PCs",
        ],
    )
    add_footer(slide)


def slide_what_is_module(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "What is a Module?", "The 'verb' of Ansible automation")
    add_bullets(
        slide,
        [
            "A standalone unit of code pushed to managed nodes",
            "Performs ONE specific, well-defined action",
            "Checks current state → compares to desired state → changes only if needed",
            "Returns structured JSON: changed, failed, msg, and more",
            "",
            "Example:",
        ],
        top=1.5,
        height=2.2,
        font_size=17,
    )
    add_code_block(
        slide,
        "- name: Ensure staging directory exists\n"
        "  ansible.windows.win_file:\n"
        "    path: C:\\Temp\\Office2016\n"
        "    state: directory",
        top=4.0,
        height=1.6,
    )
    add_footer(slide)


def slide_module_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "How Modules Execute", "Agentless push-run-return cycle")
    steps = [
        ("1", "AWX / Control Node sends module + parameters"),
        ("2", "Module copied to target Windows PC"),
        ("3", "Module runs locally on the managed node"),
        ("4", "Checks state and applies changes if needed"),
        ("5", "Returns JSON result to Ansible"),
        ("6", "Playbook continues based on result (ok / changed / failed)"),
    ]
    y = 1.6
    for num, text in steps:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.45), Inches(0.45)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        cp = circle.text_frame.paragraphs[0]
        cp.text = num
        cp.font.size = Pt(14)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER
        tb = slide.shapes.add_textbox(Inches(1.35), Inches(y), Inches(11), Inches(0.5))
        tp = tb.text_frame.paragraphs[0]
        tp.text = text
        tp.font.size = Pt(18)
        tp.font.color.rgb = DARK_TEXT
        y += 0.75
    add_footer(slide)


def slide_why_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Why We Use Modules", "Six key benefits")
    cards = [
        ("Idempotency", "Safe to re-run playbooks\nNo duplicate changes"),
        ("Declarative", "Describe WHAT you want\nModule handles HOW"),
        ("Structured Output", "JSON results for\nconditionals & registers"),
        ("Cross-Platform", "Same patterns,\nOS-specific modules"),
        ("Security", "Reviewable parameters\nvs raw shell scripts"),
        ("Ecosystem", "1000s of modules for\ncloud, network, DB, etc."),
    ]
    positions = [(0.5, 1.5), (4.6, 1.5), (8.7, 1.5), (0.5, 4.0), (4.6, 4.0), (8.7, 4.0)]
    for (title, body), (x, y) in zip(cards, positions):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = TEAL
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
    add_footer(slide)


def slide_module_vs_shell(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Module vs Command / Shell")
    rows = [
        ("Aspect", "Module", "Shell / Command"),
        ("Idempotent", "Yes ✓", "No ✗"),
        ("Declarative", "Yes ✓", "Imperative"),
        ("Output", "Structured JSON", "Raw stdout/stderr"),
        ("Best for", "Config management", "Legacy tools, no module"),
        ("Change detect", "Automatic", "Manual (changed_when)"),
    ]
    y = 1.6
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            x = 0.6 + j * 4.0
            w = 3.8 if j > 0 else 3.5
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.55)
            )
            box.fill.solid()
            if i == 0:
                box.fill.fore_color.rgb = NAVY
                color = WHITE
                bold = True
            else:
                box.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_BG
                color = DARK_TEXT
                bold = False
            box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p = box.text_frame.paragraphs[0]
            p.text = cell
            p.font.size = Pt(14 if i == 0 else 13)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER
            box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        y += 0.58
    tip = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12), Inches(0.8))
    tp = tip.text_frame.paragraphs[0]
    tp.text = "Rule: Prefer modules. Use win_shell only when no module exists (e.g., ODT setup.exe)."
    tp.font.size = Pt(14)
    tp.font.color.rgb = ACCENT
    tp.font.italic = True
    add_footer(slide)


def slide_categories(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Module Categories", "Common module families")
    add_bullets(
        slide,
        [
            "Files & Directories:  file, win_file, copy, win_copy, template, win_template",
            "Packages:             yum, apt, win_package",
            "Services:             service, win_service",
            "Information:          setup, win_stat, win_reg_stat, debug",
            "System Control:       reboot, win_reboot, user, win_user",
            "Cloud & Network:      amazon.aws.*, azure.azcollection.*, cisco.ios.*",
            "",
            "This project uses ansible.windows collection modules exclusively for Windows tasks.",
        ],
        font_size=17,
    )
    add_footer(slide)


def slide_anatomy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Anatomy of a Module Task", "FQCN + parameters + handlers")
    add_code_block(
        slide,
        "- name: Copy setup.exe from LAN share          # Task name\n"
        "  ansible.windows.win_copy:                      # FQCN\n"
        "    src: \"{{ office2016_lan_source_path }}\\\\setup.exe\"\n"
        "    dest: \"{{ office2016_staging_path }}\\\\setup.exe\"\n"
        "    remote_src: true                             # Parameters\n"
        "  register: copy_result                          # Save output\n"
        "  when: office2016_lan_setup.stat.exists         # Conditional",
        top=1.5,
        height=3.5,
    )
    add_bullets(
        slide,
        [
            "FQCN = namespace.collection.module_name (e.g., ansible.windows.win_copy)",
            "register: saves JSON output for use in when: conditions",
        ],
        top=5.3,
        height=1.5,
        font_size=15,
    )
    add_footer(slide)


def slide_idempotency(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Idempotency Explained", "Run once or a hundred times — same result")
    add_bullets(
        slide,
        [
            "Idempotent = applying config multiple times does NOT cause harm or duplicates",
            "",
            "Our Office 2016 role:",
            ("1. win_reg_stat checks if Office 16.0 registry key exists", 1),
            ("2. If installed → skip entire install block", 1),
            ("3. If not installed → copy files, run setup.exe, verify", 1),
            "",
            "First AWX run:  changed: true  (Office installed)",
            "Second AWX run: ok / skipped   (Office already present)",
            "",
            "Enables safe weekly scheduled deployments across the LAN!",
        ],
        font_size=17,
    )
    add_footer(slide)


def slide_example_reg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Example 1: win_reg_stat", "Check if Office 2016 is installed")
    add_code_block(
        slide,
        "- name: Check if Microsoft Office 2016 is already installed\n"
        "  ansible.windows.win_reg_stat:\n"
        "    path: HKLM:\\SOFTWARE\\Microsoft\\Office\\16.0\\Common\\InstallRoot\n"
        "    name: Path\n"
        "  register: office2016_reg\n"
        "  failed_when: false",
        top=1.5,
        height=2.4,
    )
    add_bullets(
        slide,
        [
            "Reads Windows registry safely on remote host",
            "Returns exists: true/false and value for install path",
            "Used with when: to skip install on already-configured PCs",
        ],
        top=4.2,
        font_size=16,
    )
    add_footer(slide)


def slide_example_file(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Example 2: win_file & win_stat", "Prepare and validate paths")
    add_code_block(
        slide,
        "- name: Ensure Office staging directory exists\n"
        "  ansible.windows.win_file:\n"
        "    path: \"{{ office2016_staging_path }}\"\n"
        "    state: directory\n\n"
        "- name: Verify LAN source path is reachable\n"
        "  ansible.windows.win_stat:\n"
        "    path: \"{{ office2016_lan_source_path }}\\\\setup.exe\"\n"
        "  register: office2016_lan_setup\n"
        "  failed_when: not office2016_lan_setup.stat.exists",
        top=1.5,
        height=3.2,
    )
    add_bullets(
        slide,
        [
            "win_file: idempotent directory creation",
            "win_stat: validates UNC share before copy — fails fast with clear error",
        ],
        top=5.0,
        font_size=16,
    )
    add_footer(slide)


def slide_example_copy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Example 3: win_copy & win_template", "Deploy Office files from LAN share")
    add_code_block(
        slide,
        "- name: Copy setup.exe from LAN share to staging\n"
        "  ansible.windows.win_copy:\n"
        "    src: \"{{ office2016_lan_source_path }}\\\\setup.exe\"\n"
        "    dest: \"{{ office2016_staging_path }}\\\\setup.exe\"\n"
        "    remote_src: true\n\n"
        "- name: Generate configuration.xml\n"
        "  ansible.windows.win_template:\n"
        "    src: configuration.xml.j2\n"
        "    dest: \"{{ office2016_staging_path }}\\\\configuration.xml\"",
        top=1.5,
        height=3.4,
    )
    add_bullets(
        slide,
        [
            "win_copy with remote_src: copies from LAN UNC path on the Windows target",
            "win_template: renders Jinja2 variables into Office ODT configuration XML",
        ],
        top=5.2,
        font_size=16,
    )
    add_footer(slide)


def slide_example_reboot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Example 4: win_reboot & debug", "Complete the deployment lifecycle")
    add_code_block(
        slide,
        "- name: Reboot when setup requires it (exit 3010)\n"
        "  ansible.windows.win_reboot:\n"
        "    reboot_timeout: 900\n"
        "  when:\n"
        "    - office2016_reboot in ['always', 'if_required']\n"
        "    - office2016_install.rc == 3010\n\n"
        "- name: Report successful installation\n"
        "  ansible.builtin.debug:\n"
        "    msg: \"Office 2016 installed at {{ office2016_verify.value }}\"",
        top=1.5,
        height=3.2,
    )
    add_bullets(
        slide,
        [
            "win_reboot: safely reboots and waits for WinRM to return",
            "debug: prints success message visible in AWX job output",
        ],
        top=5.0,
        font_size=16,
    )
    add_footer(slide)


def slide_when_shell(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "When to Use win_shell", "No module for ODT setup.exe")
    add_code_block(
        slide,
        "- name: Install Microsoft Office 2016 silently via ODT\n"
        "  ansible.windows.win_shell: |\n"
        "    $setup = Join-Path '{{ office2016_staging_path }}' 'setup.exe'\n"
        "    $config = Join-Path '{{ office2016_staging_path }}' 'configuration.xml'\n"
        "    $p = Start-Process -FilePath $setup `\n"
        "      -ArgumentList \"/configure `\"$config`\"\" -PassThru -Wait\n"
        "    exit $p.ExitCode\n"
        "  register: office2016_install\n"
        "  failed_when: office2016_install.rc not in [0, 3010]",
        top=1.5,
        height=3.5,
    )
    add_bullets(
        slide,
        [
            "Microsoft ODT has no dedicated Ansible module",
            "Shell wraps installer; failed_when treats exit 3010 (reboot needed) as success",
        ],
        top=5.3,
        font_size=16,
    )
    add_footer(slide)


def slide_best_practices(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Best Practices")
    add_bullets(
        slide,
        [
            "Prefer modules over shell commands whenever a module exists",
            "Use FQCN: ansible.builtin.debug instead of bare debug",
            "Give every task a descriptive name for AWX logs",
            "Use register: to capture output for when: conditionals",
            "Pin collection versions in requirements.yml for AWX",
            "Test idempotency — run playbook twice; second run should be mostly ok",
            "Use failed_when: for expected non-zero exit codes (e.g., 3010)",
        ],
        font_size=18,
    )
    add_footer(slide)


def slide_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Summary")
    add_bullets(
        slide,
        [
            "Module = reusable action unit (the 'verb' of Ansible)",
            "Why modules: idempotency, declarative config, structured output, security",
            "FQCN format: namespace.collection.module_name",
            "Collections extend Ansible (this project: ansible.windows)",
            "Real examples: win_reg_stat, win_file, win_copy, win_template, win_reboot",
            "Use shell only when no suitable module exists",
        ],
        font_size=19,
    )
    add_footer(slide)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5), Inches(13.333), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1))
    p = t.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    t2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.5))
    for i, line in enumerate(
        [
            "Questions?",
            "",
            "Docs: docs.ansible.com | Collection: ansible.windows",
            "Project: github.com/hussaini-8024/Project.git",
        ]
    ):
        para = t2.text_frame.paragraphs[0] if i == 0 else t2.text_frame.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_what_is_ansible(prs)
    slide_what_is_module(prs)
    slide_module_flow(prs)
    slide_why_modules(prs)
    slide_module_vs_shell(prs)
    slide_categories(prs)
    slide_anatomy(prs)
    slide_idempotency(prs)
    slide_example_reg(prs)
    slide_example_file(prs)
    slide_example_copy(prs)
    slide_example_reboot(prs)
    slide_when_shell(prs)
    slide_best_practices(prs)
    slide_summary(prs)
    slide_thank_you(prs)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
